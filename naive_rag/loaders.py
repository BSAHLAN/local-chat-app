"""Read supported files into plain text, dispatched by extension."""
from __future__ import annotations

import logging
import os

logger = logging.getLogger(__name__)

TEXT_EXTENSIONS = {
    ".txt", ".md", ".markdown", ".py", ".js", ".ts", ".java", ".c", ".cpp",
    ".h", ".go", ".rs", ".rb", ".sh", ".json", ".csv", ".tsv", ".html",
    ".htm", ".xml", ".yaml", ".yml", ".toml", ".ini", ".cfg", ".log", ".rst",
}
PDF_EXTENSIONS = {".pdf"}
DOCX_EXTENSIONS = {".docx"}
PPTX_EXTENSIONS = {".pptx"}
XLSX_EXTENSIONS = {".xlsx"}

SUPPORTED_EXTENSIONS = (
    TEXT_EXTENSIONS | PDF_EXTENSIONS | DOCX_EXTENSIONS
    | PPTX_EXTENSIONS | XLSX_EXTENSIONS
)


def _load_text(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def _load_pdf(path: str) -> str:
    from pypdf import PdfReader

    reader = PdfReader(path)
    return "\n".join((page.extract_text() or "") for page in reader.pages)


def _load_docx(path: str) -> str:
    import docx

    document = docx.Document(path)
    return "\n".join(p.text for p in document.paragraphs)


def _load_pptx(path: str) -> str:
    from pptx import Presentation

    prs = Presentation(path)
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
    return "\n".join(parts)


def _load_xlsx(path: str) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(path, read_only=True, data_only=True)
    parts: list[str] = []
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                parts.append("\t".join(cells))
    return "\n".join(parts)


def load_file(path: str) -> str | None:
    """Return extracted text for a supported file, else None.

    Raises if a supported file is corrupt; the caller decides whether to skip.
    """
    ext = os.path.splitext(path)[1].lower()
    if ext in TEXT_EXTENSIONS:
        return _load_text(path)
    if ext in PDF_EXTENSIONS:
        return _load_pdf(path)
    if ext in DOCX_EXTENSIONS:
        return _load_docx(path)
    if ext in PPTX_EXTENSIONS:
        return _load_pptx(path)
    if ext in XLSX_EXTENSIONS:
        return _load_xlsx(path)
    logger.debug("Skipping unsupported file: %s", path)
    return None
